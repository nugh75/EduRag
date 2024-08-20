import os
import logging
import time
import streamlit as st
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.docstore.document import Document
import pdfplumber
from PyPDF2 import PdfReader
from PyPDF2.errors import PdfReadError
from pdfminer.high_level import extract_text
from pdfminer.pdfparser import PDFSyntaxError
from docx import Document as DocxDocument
import odf.opendocument as odf
from odf.text import P
from pptx import Presentation

def create_database():
    logging.basicConfig(
        level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
    )

    st.write("### Crea database indicizzato da file")
    uploaded_files = st.file_uploader(
        "Carica i tuoi file:", type=["pdf", "odt", "docx", "pptx", "txt"], accept_multiple_files=True
    )

    if "metadata_list" not in st.session_state:
        st.session_state.metadata_list = []

    if "metadata_confirmed" not in st.session_state:
        st.session_state.metadata_confirmed = False

    if uploaded_files and not st.session_state.metadata_confirmed:
        st.write("### Fase 1: Inserimento e Verifica dei Metadati")

        updated_metadata = {}
        all_filled = True

        for i, uploaded_file in enumerate(uploaded_files):
            st.write(f"#### Documento: {uploaded_file.name}")

            if f"title_{i}" not in updated_metadata:
                initial_metadata = extract_metadata(uploaded_file)
                updated_metadata[f"title_{i}"] = initial_metadata["title"]
                updated_metadata[f"author_{i}"] = initial_metadata["author"]

            title_input = st.text_input(
                f"Titolo per '{uploaded_file.name}':", updated_metadata[f"title_{i}"], key=f"title_{i}"
            )
            author_input = st.text_input(
                f"Autore per '{uploaded_file.name}':", updated_metadata[f"author_{i}"], key=f"author_{i}"
            )

            if title_input.strip() == "" or author_input.strip() == "":
                all_filled = False

            updated_metadata[f"title_{i}"] = title_input
            updated_metadata[f"author_{i}"] = author_input

        if st.button("Carica Metadati"):
            if all_filled:
                st.session_state.metadata_list = [
                    {
                        "file": uploaded_file,
                        "metadata": {
                            "title": updated_metadata[f"title_{i}"],
                            "author": updated_metadata[f"author_{i}"]
                        }
                    }
                    for i, uploaded_file in enumerate(uploaded_files)
                ]
                st.session_state.metadata_confirmed = True
                st.success("Metadati caricati con successo. Procedi al passo successivo.")
            else:
                st.error("Compila tutti i campi dei metadati prima di caricarli.")

    if st.session_state.metadata_confirmed:
        st.write("### Fase 2: Embedding e Creazione dell'Indice FAISS")
        st.info("Caricamento metadati...")

        all_documents = []
        for i, entry in enumerate(st.session_state.metadata_list):
            file = entry["file"]
            metadata = entry["metadata"]

            documents, structured_content = load_file(file, metadata)
            all_documents.extend(documents)

            st.session_state.metadata_list[i]["structured_content"] = structured_content

        chunk_size = 1000
        chunk_overlap = 20
        min_chunk_length = 50
        embeddings = HuggingFaceEmbeddings(
            model_name="sentence-transformers/all-MiniLM-L12-v2"
        )
        subfolder_name = st.text_input(
            "Nome del database indicizzato:", "Inserisci il nome del database"
        )

        index_description = st.text_area(
            "Descrizione del database:",
            "Inserisci una breve descrizione del database indicizzato",
        )

        faiss_index_folder = os.path.join("app/db", subfolder_name)

        if not subfolder_name:
            st.error("Inserisci un nome valido per il database indicizzato.")
            return

        if st.button("Procedi con l'Embedding e la Creazione dell'Indice"):
            progress_text = st.empty()
            progress_text.text("Divisione dei documenti in chunk...")

            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size, chunk_overlap=chunk_overlap
            )

            chunked_documents = []
            for entry in st.session_state.metadata_list:
                title = entry["metadata"]["title"]
                author = entry["metadata"]["author"]
                structured_content = entry["structured_content"]

                updated_metadata = {
                    "title": title,
                    "author": author,
                }
                for page_number, page_content in enumerate(structured_content, start=1):
                    chunks = text_splitter.split_text(page_content)
                    for chunk in chunks:
                        chunk_metadata = updated_metadata.copy()
                        chunk_metadata["page_number"] = page_number

                        chunked_documents.append(
                            Document(page_content=chunk, metadata=chunk_metadata)
                        )

            splits = [sp for sp in chunked_documents if len(sp.page_content) >= min_chunk_length]

            if not splits:
                st.error("Nessun chunk valido è stato creato. Verifica i parametri di suddivisione.")
                logging.error("Nessun chunk valido è stato creato. Verifica i parametri di suddivisione.")
                return

            st.header("Primi 3 Chunk Estratti")
            for i, chunk in enumerate(splits[:3]):
                st.subheader(f"Chunk {i + 1} - Metadati:")
                st.json(chunk.metadata)

                st.subheader(f"Contenuto del Chunk {i + 1}:")
                st.write(chunk.page_content)

            progress_text.text("Creazione dell'indice FAISS...")

            for _ in range(5):
                for dots in range(1, 6):
                    progress_text.text(f"Sto elaborando il database{'.' * dots}")
                    time.sleep(0.5)

            docs = [
                Document(page_content=split.page_content, metadata=split.metadata)
                for split in splits
            ]

            if not docs:
                st.error(
                    "Nessun documento da indicizzare. Verifica che i documenti contengano testo valido."
                )
                logging.error(
                    "Nessun documento da indicizzare. Verifica che i documenti contengano testo valido."
                )
                return

            index = FAISS.from_documents(docs, embeddings)

            if not os.path.exists(faiss_index_folder):
                os.makedirs(faiss_index_folder)

            index.save_local(faiss_index_folder)

            description_file_path = os.path.join(faiss_index_folder, "description.txt")
            with open(description_file_path, "w") as desc_file:
                desc_file.write(f"Descrizione dell'indice: {index_description}\n")
                desc_file.write("Titoli dei documenti:\n")
                for entry in st.session_state.metadata_list:
                    desc_file.write(f"- {entry['metadata']['title']}\n")

            progress_text.text(f"Indice '{subfolder_name}' creato con successo.")


def extract_metadata(file):
    metadata = {"title": "Sconosciuto", "author": "Sconosciuto"}
    try:
        if file.name.endswith(".pdf"):
            reader = PdfReader(file)
            doc_info = reader.metadata
            metadata["title"] = doc_info.get("/Title", "Sconosciuto")
            metadata["author"] = doc_info.get("/Author", "Sconosciuto")
        elif file.name.endswith(".docx"):
            doc = DocxDocument(file)
            metadata["title"] = doc.core_properties.title or "Sconosciuto"
            metadata["author"] = doc.core_properties.author or "Sconosciuto"
        elif file.name.endswith(".odt"):
            odt_file = odf.load(file)
            metadata["title"] = odt_file.meta.get("title", "Sconosciuto")
            metadata["author"] = odt_file.meta.get("creator", "Sconosciuto")
        elif file.name.endswith(".pptx"):
            ppt = Presentation(file)
            metadata["title"] = ppt.core_properties.title or "Sconosciuto"
            metadata["author"] = ppt.core_properties.author or "Sconosciuto"
        elif file.name.endswith(".txt"):
            metadata["title"] = file.name
            metadata["author"] = "Sconosciuto"
    except Exception as e:
        logging.error(f"Errore nell'estrazione dei metadati dal file {file.name}: {e}")

    return metadata


def load_file(file, metadata):
    try:
        if file.name.endswith(".pdf"):
            structured_content = extract_structured_content_pdf(file)
        elif file.name.endswith(".docx"):
            structured_content = extract_structured_content_docx(file)
        elif file.name.endswith(".odt"):
            structured_content = extract_structured_content_odt(file)
        elif file.name.endswith(".pptx"):
            structured_content = extract_structured_content_pptx(file)
        elif file.name.endswith(".txt"):
            structured_content = extract_structured_content_txt(file)
        documents = create_documents_with_metadata(structured_content, metadata)
        return documents, structured_content
    except Exception as e:
        logging.error(f"Errore nel caricamento del file {file.name}: {str(e)}")
        return [], []


def extract_structured_content_pdf(file):
    structured_content = []
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            structured_content.append(text)
    return structured_content


def extract_structured_content_docx(file):
    doc = DocxDocument(file)
    return [para.text for para in doc.paragraphs if para.text.strip()]


def extract_structured_content_odt(file):
    odt_file = odf.load(file)
    text_content = []
    for elem in odt_file.getElementsByType(P):
        text_content.append(str(elem))
    return text_content


def extract_structured_content_pptx(file):
    ppt = Presentation(file)
    text_content = []
    for slide in ppt.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text_content.append("\n".join(slide_text))
    return text_content


def extract_structured_content_txt(file):
    text = file.read().decode("utf-8")
    return text.splitlines()


def create_documents_with_metadata(chunks, metadata):
    docs = []
    for chunk in chunks:
        doc = Document(page_content=chunk, metadata=metadata)
        docs.append(doc)
    return docs


if __name__ == "__main__":
    create_database()
